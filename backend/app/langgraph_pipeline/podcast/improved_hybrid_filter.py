"""
Improved Hybrid Filter V2 (get_images() 방식)
==============================================

핵심 변경사항:
- get_text('dict') → get_images() + get_image_bbox()
- 모든 이미지 감지 (배경 레이어 포함)
- 만화 콘텐츠 정상 인식
"""

import os
import vertexai
import textwrap
import json
from dataclasses import dataclass
from typing import List, Dict
from pptx import Presentation
from vertexai.generative_models import GenerativeModel, Part

# [1] 인증 설정
SERVICE_ACCOUNT_FILE = "vertex-ai-service-account.json"
os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = SERVICE_ACCOUNT_FILE
PROJECT_ID = "alan-document-lab" 
vertexai.init(project=PROJECT_ID, location="us-central1")

# [2] Gemini 2.5 Flash 모델 로드
model = GenerativeModel("gemini-2.5-flash")

@dataclass
class ImageMetadata:
    image_id: str
    slide_number: int
    area_percentage: float
    left: float
    top: float
    adjacent_text: str
    slide_title: str
    image_bytes: bytes = None
    is_core_content: bool = False
    filter_reason: str = ""

# 1. 통합 이미지 추출기 (PPTX + PDF 지원)
class UniversalImageExtractor:
    """
    모든 형식에서 이미지 메타데이터 추출
    V2: get_images() 방식으로 모든 이미지 감지
    """
    
    def extract(self, file_path: str) -> List[ImageMetadata]:
        from pathlib import Path
        
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pptx':
            return self._extract_from_pptx(file_path)
        elif ext == '.pdf':
            return self._extract_from_pdf_v2(file_path)
        else:
            raise ValueError(f"지원하지 않는 형식: {ext}")
    
    def _extract_from_pptx(self, pptx_path: str) -> List[ImageMetadata]:
        """PPTX에서 이미지 추출 (기존 방식)"""
        if not os.path.exists(pptx_path):
            return []
        
        prs = Presentation(pptx_path)
        metadata_list = []
        slide_width, slide_height = prs.slide_width.inches, prs.slide_height.inches
        slide_area = slide_width * slide_height

        for s_idx, slide in enumerate(prs.slides, 1):
            slide_title = slide.shapes.title.text if slide.shapes.title else "No Title"
            all_text = " ".join([s.text for s in slide.shapes if hasattr(s, "text")])
            
            img_idx = 1
            for shape in slide.shapes:
                if shape.shape_type == 13 or hasattr(shape, 'image'):
                    w, h = shape.width.inches, shape.height.inches
                    area_pct = ((w * h) / slide_area) * 100
                    metadata_list.append(ImageMetadata(
                        image_id=f"S{s_idx:02d}_IMG{img_idx:03d}",
                        slide_number=s_idx,
                        area_percentage=area_pct,
                        left=shape.left.inches,
                        top=shape.top.inches,
                        adjacent_text=all_text.replace('\n', ' ').strip(),
                        slide_title=slide_title,
                        image_bytes=shape.image.blob
                    ))
                    img_idx += 1
        
        return metadata_list
    
    def _extract_text_with_ocr(self, page, min_length: int = 100) -> str:
        """페이지에서 텍스트 추출 (필요시 OCR)"""
        text = page.get_text()
        text_length = len(text.strip())
        
        if text_length >= min_length:
            return text
        
        try:
            from paddleocr import PaddleOCR
            
            if not hasattr(self, '_ocr_engine'):
                os.environ['FLAGS_log_level'] = '3'
                os.environ['PPOCR_SHOW_LOG'] = 'False'
                
                print(f"      → PaddleOCR 초기화 중...")
                self._ocr_engine = PaddleOCR(lang='korean', use_textline_orientation=True)
            
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            
            import numpy as np
            from PIL import Image
            from io import BytesIO
            
            img = Image.open(BytesIO(img_data))
            img_array = np.array(img)
            
            result = self._ocr_engine.ocr(img_array)
            
            if result and result[0]:
                lines = []
                for line in result[0]:
                    if line and len(line) >= 2:
                        ocr_text = line[1][0]
                        lines.append(ocr_text)
                
                ocr_result = "\n".join(lines)
                print(f"      → 페이지 OCR: {text_length}자 → {len(ocr_result)}자")
                return ocr_result if ocr_result else text
        
        except ImportError:
            pass
        except Exception as e:
            print(f"      ⚠️  OCR 실패: {e}")
        
        return text
    
    def _extract_page_title(self, page_text: str) -> str:
        """페이지 제목 추출"""
        lines = page_text.strip().split('\n')
        for line in lines:
            line = line.strip()
            if len(line) > 3 and not line.startswith('☞'):
                return line[:50]
        return "페이지 제목 없음"
    
    def _extract_from_pdf_v2(self, pdf_path: str) -> List[ImageMetadata]:
        """
        PDF에서 이미지 추출 (V2: get_images() 방식)
        
        핵심 변경:
        - get_text('dict') → get_images() + get_image_bbox()
        - 모든 이미지 감지 (배경 레이어 포함)
        """
        try:
            import fitz
        except ImportError:
            print("   ❌ PyMuPDF가 설치되지 않았습니다.")
            return []
        
        if not os.path.exists(pdf_path):
            return []
        
        metadata_list = []
        
        # 필터링 기준
        MIN_WIDTH = 40
        MIN_HEIGHT = 40
        MIN_AREA_PCT = 3.0      # 3% 미만: 레이블/아이콘
        MAX_AREA_PCT = 90.0     # 90% 이상: 배경
        MIN_PIXEL_AREA = 1000
        MAX_ASPECT_RATIO = 6.0  # 6:1 이상: 제목/텍스트
        
        total_images = 0
        filtered_background = 0
        filtered_aspect = 0
        filtered_area = 0
        filtered_size = 0
        
        try:
            doc = fitz.open(pdf_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # 페이지 정보
                page_area = page.rect.width * page.rect.height
                page_text = self._extract_text_with_ocr(page, min_length=100)
                page_title = self._extract_page_title(page_text)
                
                # ===== get_images()로 모든 이미지 감지 =====
                images = page.get_images(full=True)
                total_images += len(images)
                
                print(f"      [P{page_num+1}] 총 {len(images)}개 이미지 발견")
                
                for img in images:
                    try:
                        xref = img[0]
                        
                        # bbox 가져오기
                        try:
                            bbox = page.get_image_bbox(img)
                        except:
                            continue
                        
                        if not bbox or bbox.is_empty or bbox.is_infinite:
                            continue
                        
                        x0, y0, x1, y1 = bbox
                        width = x1 - x0
                        height = y1 - y0
                        area_pct = (width * height) / page_area * 100
                        
                        debug_msg = f"      [P{page_num+1}] {area_pct:.1f}%"
                        
                        # ===== 필터 1: 배경 제외 (90% 이상) =====
                        if area_pct > MAX_AREA_PCT:
                            filtered_background += 1
                            print(debug_msg + f" → 배경 제외 ❌")
                            continue
                        
                        # ===== 필터 2: 가로세로비 =====
                        if width > 0 and height > 0:
                            aspect_ratio = max(width, height) / min(width, height)
                            if aspect_ratio > MAX_ASPECT_RATIO:
                                filtered_aspect += 1
                                print(debug_msg + f" → 가로세로비 제외 ({aspect_ratio:.1f}:1) ❌")
                                continue
                        
                        # ===== 필터 3: 작은 면적 =====
                        pixel_area = width * height
                        if pixel_area < MIN_PIXEL_AREA:
                            filtered_area += 1
                            print(debug_msg + f" → 작은 면적 제외 ❌")
                            continue
                        
                        # ===== 필터 4: 절대 크기 =====
                        if width < MIN_WIDTH or height < MIN_HEIGHT:
                            filtered_size += 1
                            print(debug_msg + f" → 작은 크기 제외 ❌")
                            continue
                        
                        # ===== 필터 5: 상대 크기 =====
                        if area_pct < MIN_AREA_PCT:
                            filtered_size += 1
                            print(debug_msg + f" → 상대 크기 제외 ({area_pct:.1f}%) ❌")
                            continue
                        
                        # ===== 통과! =====
                        print(debug_msg + " → 최종 추출 ✅✅✅")
                        
                        # 이미지 추출
                        try:
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                        except:
                            pix = page.get_pixmap(clip=fitz.Rect(bbox), dpi=150)
                            image_bytes = pix.tobytes("png")
                        
                        metadata_list.append(ImageMetadata(
                            image_id=f"P{page_num+1:02d}_IMG{len(metadata_list)+1:03d}",
                            slide_number=page_num + 1,
                            area_percentage=area_pct,
                            left=x0,
                            top=y0,
                            adjacent_text=page_text.replace('\n', ' ').strip(),
                            slide_title=page_title,
                            image_bytes=image_bytes
                        ))
                    
                    except Exception as e:
                        print(f"      ⚠️ 이미지 처리 실패: {e}")
                        continue
            
            doc.close()
        
        except Exception as e:
            print(f"   ❌ PDF 처리 실패: {e}")
            return []
        
        # 통계
        print(f"\n   📊 PDF 이미지 분석:")
        print(f"      - 전체 이미지: {total_images}개")
        print(f"   🔍 필터링 통계:")
        print(f"      - 배경 제외: {filtered_background}개")
        print(f"      - 가로세로비: {filtered_aspect}개")
        print(f"      - 작은 면적: {filtered_area}개")
        print(f"      - 작은 크기: {filtered_size}개")
        print(f"   ✅ 최종 추출: {len(metadata_list)}개 이미지\n")
        
        return metadata_list


# 2. 개선된 하이브리드 필터 파이프라인
class ImprovedHybridFilterPipeline:
    def __init__(self, auto_extract_keywords: bool = True):
        self.auto_extract = auto_extract_keywords
        
        self.UNIVERSAL_PATTERNS = [
            '학습', '활동', '문제', '예제', '연습',
            '생각', '알아보', '살펴보', '정리',
            '목표', '개념', '원리', '법칙', '정의',
            '단원', '차시',
            '그림', '도표', '표', '차트', '그래프',
            '예시', '사례', '모형', '구조'
        ]
        
        self.DECORATION_PATTERNS = [
            '로고', 'logo', '출처', '참고', '아이콘', 'icon'
        ]
        
        self.document_keywords = []

    def extract_keywords_from_document(self, file_path: str):
        """문서에서 자동으로 키워드 추출"""
        if not self.auto_extract:
            return
        
        from pathlib import Path
        
        print("📚 문서 분석하여 키워드 자동 추출 중...")
        
        ext = Path(file_path).suffix.lower()
        all_text = []
        
        if ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
        
        elif ext == '.pdf':
            import pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            all_text.append(text)
            except Exception as e:
                print(f"   ⚠️ PDF 텍스트 추출 실패, 범용 패턴만 사용")
                return
        
        else:
            print(f"   ⚠️ 지원하지 않는 형식: {ext}")
            return
        
        full_text = "\n".join(all_text)[:5000]
        
        prompt = f"""
다음 강의 자료에서 **핵심 키워드 20개**를 추출하세요.

# 문서 내용
{full_text}

# 조건
- 개념어, 전문 용어, 주제어만 포함
- JSON 형식: {{"keywords": ["키워드1", "키워드2", ...]}}
"""
        
        try:
            response = model.generate_content(prompt)

            # ✅ 토큰 추출
            usage = response.usage_metadata
            in_t = usage.prompt_token_count
            out_t = usage.candidates_token_count
            cost = (in_t / 1_000_000 * 0.075) + (out_t / 1_000_000 * 0.30)
            
            print(f"📊 [키워드 추출] 토큰: {usage.total_token_count:,} (In: {in_t}/Out: {out_t}) / 비용: ${cost:.6f}")

            text = response.text.strip()
            
            if "```json" in text:
                text = text.split("```json")[1].split("```")[0].strip()
            elif "```" in text:
                text = text.split("```")[1].split("```")[0].strip()
            
            data = json.loads(text)
            self.document_keywords = data.get("keywords", [])
            
            print(f"   ✅ 추출된 키워드: {', '.join(self.document_keywords[:10])}")
        
        except Exception as e:
            print(f"   ⚠️ 자동 추출 실패, 범용 패턴만 사용")
            self.document_keywords = []

    def step1_rule_check(self, meta: ImageMetadata):
        """규칙 기반 1차 필터"""
        context = f"{meta.slide_title} {meta.adjacent_text}".lower()
        
        has_deco = any(kw in context for kw in self.DECORATION_PATTERNS)
        is_corner = (meta.left < 1.0 and meta.top < 1.0) or (meta.left > 8.0 and meta.top < 1.0)
        
        if is_corner and meta.area_percentage < 5.0 and not any(kw in context for kw in self.UNIVERSAL_PATTERNS):
            return "EXCLUDE", "Static Decoration (Corner)"
        
        if has_deco and meta.area_percentage < 8.0:
            return "EXCLUDE", "Decorative element"
        
        has_universal = any(p in context for p in self.UNIVERSAL_PATTERNS)
        has_document_kw = any(kw in context for kw in self.document_keywords)
        
        if meta.area_percentage > 15.0 and (has_universal or has_document_kw):
            return "INCLUDE", f"Core content ({meta.area_percentage:.1f}% + pattern)"
        
        if has_document_kw and meta.area_percentage > 10.0:
            matched = [kw for kw in self.document_keywords if kw in context]
            return "INCLUDE", f"Document keyword: {', '.join(matched[:2])}"
        
        return "PENDING", "Requires AI Vision Check"

    def step2_gemini_check(self, meta: ImageMetadata, max_retries=3):
        """AI Vision으로 2차 판단"""
        import time
        
        for attempt in range(max_retries):
            try:
                image_part = Part.from_data(data=meta.image_bytes, mime_type="image/png")
                
                keyword_list = ', '.join(list(self.document_keywords)[:15]) if self.document_keywords else "일반 학습 내용"
                
                prompt = f"""
이 강의의 핵심 주제: {keyword_list}

이 이미지가 위 주제들과 관련있는지 판단하세요.

주변 텍스트: "{meta.adjacent_text}"

판단:
- 학습에 필요한 핵심 자료 → KEEP + 이유
- 장식/로고/배경 → DISCARD + 이유

출력: KEEP 또는 DISCARD로 시작
"""
                response = model.generate_content([image_part, prompt])

                # ✅ 토큰 및 비용 계산 추가
                # ✅ Gemini 2.5 Flash 공식 단가 적용
                input_tokens = response.usage_metadata.prompt_token_count
                output_tokens = response.usage_metadata.candidates_token_count
                total_tokens = response.usage_metadata.total_token_count
                cost = (input_tokens / 1_000_000 * 0.075) + (output_tokens / 1_000_000 * 0.30)

                return response.text.strip(), total_tokens, cost
                
            except Exception as e:
                error_msg = str(e)
                
                if "429" in error_msg or "Resource exhausted" in error_msg:
                    if attempt < max_retries - 1:
                        wait_time = (attempt + 1) * 3
                        print(f"      ⚠️  Rate Limit, {wait_time}초 대기...")
                        time.sleep(wait_time)
                        continue
                    else:
                        return "DISCARD: API rate limit exceeded", 0, 0.0
                else:
                    return f"ERROR: {error_msg}", 0, 0.0
        
        return "DISCARD: Failed after all retries", 0, 0.0

    def run(self, source_path: str):
        """이미지 필터링 실행"""
        from pathlib import Path
        
        file_ext = Path(source_path).suffix.lower()
        print(f"\n🔍 분석 시작: {os.path.basename(source_path)} ({file_ext})")
        
        if self.auto_extract:
            self.extract_keywords_from_document(source_path)
        
        extractor = UniversalImageExtractor()
        all_meta = extractor.extract(source_path)
        
        print("\n" + "="*120)
        print(f"{'Slide':<6} | {'Size':<6} | {'Filter':<12} | {'Result':<12} | {'Reason'}")
        print("-" * 120)

        final_core = []
        stats = {
            'total': len(all_meta),
            'rule_pass': 0,
            'rule_drop': 0,
            'ai_keep': 0,
            'ai_drop': 0,
        }
        
        for meta in all_meta:
            decision_type, s1_reason = self.step1_rule_check(meta)

            final_status = ""
            filter_stage = ""
            detail_reason = ""

            if decision_type == "INCLUDE":
                meta.is_core_content = True
                filter_stage = "1차 (Rule)"
                final_status = "✅ PASS"
                detail_reason = s1_reason
                final_core.append(meta)
                stats['rule_pass'] += 1
                
            elif decision_type == "PENDING":
                filter_stage = "2차 (AI)"
                ai_res, tokens, cost = self.step2_gemini_check(meta)
                
                if ai_res.upper().startswith("KEEP"):
                    meta.is_core_content = True
                    final_status = "✅ KEEP"
                    stats['ai_keep'] += 1
                    final_core.append(meta)
                else:
                    final_status = "❌ DROP"
                    stats['ai_drop'] += 1
                    
                detail_reason = ai_res.replace('\n', ' ')
                
            else:
                filter_stage = "1차 (Rule)"
                final_status = "❌ DROP"
                detail_reason = s1_reason
                stats['rule_drop'] += 1

            wrapped_reason = textwrap.wrap(detail_reason, width=70)
            print(f"{meta.slide_number:<6} | {meta.area_percentage:>5.1f}% | {filter_stage:<12} | {final_status:<12} | {wrapped_reason[0]}")
            for line in wrapped_reason[1:]:
                print(f"{'':<6} | {'':<6} | {'':<12} | {'':<12} | {line}")
            print("-" * 120)

        print("\n" + "="*120)
        print("📊 최종 결과")
        print("="*120)
        
        print(f"\n총 이미지: {stats['total']}개")
        print(f"\n[1차 필터 - 규칙 기반]")
        print(f"  ✅ 통과: {stats['rule_pass']}개")
        print(f"  ❌ 제외: {stats['rule_drop']}개")
        print(f"  ⚠️  2차 이동: {stats['ai_keep'] + stats['ai_drop']}개")
        
        print(f"\n[2차 필터 - AI 판단]")
        print(f"  ✅ 통과: {stats['ai_keep']}개")
        print(f"  ❌ 제외: {stats['ai_drop']}개")
        
        total_keep = stats['rule_pass'] + stats['ai_keep']
        total_drop = stats['rule_drop'] + stats['ai_drop']
        
        print(f"\n{'='*120}")
        print(f"💎 최종 핵심 이미지: {total_keep}개 (1차: {stats['rule_pass']}개 + 2차: {stats['ai_keep']}개)")
        print(f"🗑️  제외된 이미지: {total_drop}개")
        if stats['total'] > 0:
            print(f"💰 Vision API 사용: {stats['ai_keep'] + stats['ai_drop']}회 ({(stats['ai_keep'] + stats['ai_drop'])/stats['total']*100:.1f}%)")
        print(f"{'='*120}\n")
        
        return final_core


if __name__ == "__main__":
    import sys
    
    print("\n" + "="*120)
    print("🎯 Improved Hybrid Filter V2 - 이미지 필터링")
    print("="*120)
    
    if len(sys.argv) > 1:
        source_file = sys.argv[1]
        
        if not os.path.exists(source_file):
            print(f"\n❌ 파일을 찾을 수 없습니다: {source_file}")
            sys.exit(1)
        
        auto_extract = True
        if len(sys.argv) > 2 and sys.argv[2] in ['--no-auto', '-n']:
            auto_extract = False
            print("\n⚠️  자동 키워드 추출 비활성화")
        else:
            print("\n✅ 자동 키워드 추출 활성화")
        
        try:
            pipeline = ImprovedHybridFilterPipeline(auto_extract_keywords=auto_extract)
            core_images = pipeline.run(source_file)
            
            print(f"\n{'='*120}")
            print(f"✅ 완료! 핵심 이미지: {len(core_images)}개")
            print(f"{'='*120}\n")
            
        except Exception as e:
            print(f"\n❌ 에러 발생: {e}")
            import traceback
            traceback.print_exc()
            sys.exit(1)
    
    else:
        print("\n사용법:")
        print("  python improved_hybrid_filter_v2.py <파일경로>")
        print("\n예시:")
        print("  python improved_hybrid_filter_v2.py 중등국어1.pdf")
        print("\n✅ V2 개선사항:")
        print("  - get_images() 방식으로 모든 이미지 감지")
        print("  - 만화 콘텐츠 정상 인식")
        print("  - 배경 이미지 자동 제외")
        print("="*120 + "\n")