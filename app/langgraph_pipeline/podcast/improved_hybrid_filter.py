"""
Improved Hybrid Filter V2 (get_images() ë°©ì‹)
==============================================

í•µì‹¬ ë³€ê²½ì‚¬í•­:
- get_text('dict') â†’ get_images() + get_image_bbox()
- ëª¨ë“  ì´ë¯¸ì§€ ê°ì§€ (ë°°ê²½ ë ˆì´ì–´ í¬í•¨)
- ë§Œí™” ì½˜í…ì¸  ì •ìƒ ì¸ì‹
"""

import os
import vertexai
import textwrap
import json
from dataclasses import dataclass
from typing import List, Dict
from pptx import Presentation
from vertexai.generative_models import GenerativeModel, Part

# [1] ì¸ì¦ ì„¤ì •
SERVICE_ACCOUNT_FILE = "vertex-ai-service-account.json"
os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = SERVICE_ACCOUNT_FILE
PROJECT_ID = "alan-document-lab" 
vertexai.init(project=PROJECT_ID, location="us-central1")

# [2] Gemini 2.5 Flash ëª¨ë¸ ë¡œë“œ
model = GenerativeModel("gemini-2.5-flash")

@dataclass
class ImageMetadata:
    image_id: str
    slide_number: int
    area_percentage: float
    left: float
    top: float
    adjacent_text: str
    slide_title: str
    image_bytes: bytes = None
    is_core_content: bool = False
    filter_reason: str = ""

# 1. í†µí•© ì´ë¯¸ì§€ ì¶”ì¶œê¸° (PPTX + PDF ì§€ì›)
class UniversalImageExtractor:
    """
    ëª¨ë“  í˜•ì‹ì—ì„œ ì´ë¯¸ì§€ ë©”íƒ€ë°ì´í„° ì¶”ì¶œ
    V2: get_images() ë°©ì‹ìœ¼ë¡œ ëª¨ë“  ì´ë¯¸ì§€ ê°ì§€
    """
    
    def extract(self, file_path: str) -> List[ImageMetadata]:
        from pathlib import Path
        
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pptx':
            return self._extract_from_pptx(file_path)
        elif ext == '.pdf':
            return self._extract_from_pdf_v2(file_path)
        else:
            raise ValueError(f"ì§€ì›í•˜ì§€ ì•ŠëŠ” í˜•ì‹: {ext}")
    
    def _extract_from_pptx(self, pptx_path: str) -> List[ImageMetadata]:
        """PPTXì—ì„œ ì´ë¯¸ì§€ ì¶”ì¶œ (ê¸°ì¡´ ë°©ì‹)"""
        if not os.path.exists(pptx_path):
            return []
        
        prs = Presentation(pptx_path)
        metadata_list = []
        slide_width, slide_height = prs.slide_width.inches, prs.slide_height.inches
        slide_area = slide_width * slide_height

        for s_idx, slide in enumerate(prs.slides, 1):
            slide_title = slide.shapes.title.text if slide.shapes.title else "No Title"
            all_text = " ".join([s.text for s in slide.shapes if hasattr(s, "text")])
            
            img_idx = 1
            for shape in slide.shapes:
                if shape.shape_type == 13 or hasattr(shape, 'image'):
                    w, h = shape.width.inches, shape.height.inches
                    area_pct = ((w * h) / slide_area) * 100
                    metadata_list.append(ImageMetadata(
                        image_id=f"S{s_idx:02d}_IMG{img_idx:03d}",
                        slide_number=s_idx,
                        area_percentage=area_pct,
                        left=shape.left.inches,
                        top=shape.top.inches,
                        adjacent_text=all_text.replace('\n', ' ').strip(),
                        slide_title=slide_title,
                        image_bytes=shape.image.blob
                    ))
                    img_idx += 1
        
        return metadata_list
    
    def _extract_text_with_ocr(self, page, min_length: int = 100) -> str:
        """í˜ì´ì§€ì—ì„œ í…ìŠ¤íŠ¸ ì¶”ì¶œ (í•„ìš”ì‹œ OCR)"""
        text = page.get_text()
        text_length = len(text.strip())
        
        if text_length >= min_length:
            return text
        
        try:
            from paddleocr import PaddleOCR
            
            if not hasattr(self, '_ocr_engine'):
                os.environ['FLAGS_log_level'] = '3'
                os.environ['PPOCR_SHOW_LOG'] = 'False'
                
                print(f"      â†’ PaddleOCR ì´ˆê¸°í™” ì¤‘...")
                self._ocr_engine = PaddleOCR(lang='korean', use_textline_orientation=True)
            
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            
            import numpy as np
            from PIL import Image
            from io import BytesIO
            
            img = Image.open(BytesIO(img_data))
            img_array = np.array(img)
            
            result = self._ocr_engine.ocr(img_array)
            
            if result and result[0]:
                lines = []
                for line in result[0]:
                    if line and len(line) >= 2:
                        ocr_text = line[1][0]
                        lines.append(ocr_text)
                
                ocr_result = "\n".join(lines)
                print(f"      â†’ í˜ì´ì§€ OCR: {text_length}ì â†’ {len(ocr_result)}ì")
                return ocr_result if ocr_result else text
        
        except ImportError:
            pass
        except Exception as e:
            print(f"      âš ï¸  OCR ì‹¤íŒ¨: {e}")
        
        return text
    
    def _extract_page_title(self, page_text: str) -> str:
        """í˜ì´ì§€ ì œëª© ì¶”ì¶œ"""
        lines = page_text.strip().split('\n')
        for line in lines:
            line = line.strip()
            if len(line) > 3 and not line.startswith('â˜'):
                return line[:50]
        return "í˜ì´ì§€ ì œëª© ì—†ìŒ"
    
    def _extract_from_pdf_v2(self, pdf_path: str) -> List[ImageMetadata]:
        """
        PDFì—ì„œ ì´ë¯¸ì§€ ì¶”ì¶œ (V2: get_images() ë°©ì‹)
        
        í•µì‹¬ ë³€ê²½:
        - get_text('dict') â†’ get_images() + get_image_bbox()
        - ëª¨ë“  ì´ë¯¸ì§€ ê°ì§€ (ë°°ê²½ ë ˆì´ì–´ í¬í•¨)
        """
        try:
            import fitz
        except ImportError:
            print("   âŒ PyMuPDFê°€ ì„¤ì¹˜ë˜ì§€ ì•Šì•˜ìŠµë‹ˆë‹¤.")
            return []
        
        if not os.path.exists(pdf_path):
            return []
        
        metadata_list = []
        
        # í•„í„°ë§ ê¸°ì¤€
        MIN_WIDTH = 40
        MIN_HEIGHT = 40
        MIN_AREA_PCT = 3.0      # 3% ë¯¸ë§Œ: ë ˆì´ë¸”/ì•„ì´ì½˜
        MAX_AREA_PCT = 90.0     # 90% ì´ìƒ: ë°°ê²½
        MIN_PIXEL_AREA = 1000
        MAX_ASPECT_RATIO = 6.0  # 6:1 ì´ìƒ: ì œëª©/í…ìŠ¤íŠ¸
        
        total_images = 0
        filtered_background = 0
        filtered_aspect = 0
        filtered_area = 0
        filtered_size = 0
        
        try:
            doc = fitz.open(pdf_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # í˜ì´ì§€ ì •ë³´
                page_area = page.rect.width * page.rect.height
                page_text = self._extract_text_with_ocr(page, min_length=100)
                page_title = self._extract_page_title(page_text)
                
                # ===== get_images()ë¡œ ëª¨ë“  ì´ë¯¸ì§€ ê°ì§€ =====
                images = page.get_images(full=True)
                total_images += len(images)
                
                print(f"      [P{page_num+1}] ì´ {len(images)}ê°œ ì´ë¯¸ì§€ ë°œê²¬")
                
                for img in images:
                    try:
                        xref = img[0]
                        
                        # bbox ê°€ì ¸ì˜¤ê¸°
                        try:
                            bbox = page.get_image_bbox(img)
                        except:
                            continue
                        
                        if not bbox or bbox.is_empty or bbox.is_infinite:
                            continue
                        
                        x0, y0, x1, y1 = bbox
                        width = x1 - x0
                        height = y1 - y0
                        area_pct = (width * height) / page_area * 100
                        
                        debug_msg = f"      [P{page_num+1}] {area_pct:.1f}%"
                        
                        # ===== í•„í„° 1: ë°°ê²½ ì œì™¸ (90% ì´ìƒ) =====
                        if area_pct > MAX_AREA_PCT:
                            filtered_background += 1
                            print(debug_msg + f" â†’ ë°°ê²½ ì œì™¸ âŒ")
                            continue
                        
                        # ===== í•„í„° 2: ê°€ë¡œì„¸ë¡œë¹„ =====
                        if width > 0 and height > 0:
                            aspect_ratio = max(width, height) / min(width, height)
                            if aspect_ratio > MAX_ASPECT_RATIO:
                                filtered_aspect += 1
                                print(debug_msg + f" â†’ ê°€ë¡œì„¸ë¡œë¹„ ì œì™¸ ({aspect_ratio:.1f}:1) âŒ")
                                continue
                        
                        # ===== í•„í„° 3: ì‘ì€ ë©´ì  =====
                        pixel_area = width * height
                        if pixel_area < MIN_PIXEL_AREA:
                            filtered_area += 1
                            print(debug_msg + f" â†’ ì‘ì€ ë©´ì  ì œì™¸ âŒ")
                            continue
                        
                        # ===== í•„í„° 4: ì ˆëŒ€ í¬ê¸° =====
                        if width < MIN_WIDTH or height < MIN_HEIGHT:
                            filtered_size += 1
                            print(debug_msg + f" â†’ ì‘ì€ í¬ê¸° ì œì™¸ âŒ")
                            continue
                        
                        # ===== í•„í„° 5: ìƒëŒ€ í¬ê¸° =====
                        if area_pct < MIN_AREA_PCT:
                            filtered_size += 1
                            print(debug_msg + f" â†’ ìƒëŒ€ í¬ê¸° ì œì™¸ ({area_pct:.1f}%) âŒ")
                            continue
                        
                        # ===== í†µê³¼! =====
                        print(debug_msg + " â†’ ìµœì¢… ì¶”ì¶œ âœ…âœ…âœ…")
                        
                        # ì´ë¯¸ì§€ ì¶”ì¶œ
                        try:
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                        except:
                            pix = page.get_pixmap(clip=fitz.Rect(bbox), dpi=150)
                            image_bytes = pix.tobytes("png")
                        
                        metadata_list.append(ImageMetadata(
                            image_id=f"P{page_num+1:02d}_IMG{len(metadata_list)+1:03d}",
                            slide_number=page_num + 1,
                            area_percentage=area_pct,
                            left=x0,
                            top=y0,
                            adjacent_text=page_text.replace('\n', ' ').strip(),
                            slide_title=page_title,
                            image_bytes=image_bytes
                        ))
                    
                    except Exception as e:
                        print(f"      âš ï¸ ì´ë¯¸ì§€ ì²˜ë¦¬ ì‹¤íŒ¨: {e}")
                        continue
            
            doc.close()
        
        except Exception as e:
            print(f"   âŒ PDF ì²˜ë¦¬ ì‹¤íŒ¨: {e}")
            return []
        
        # í†µê³„
        print(f"\n   ğŸ“Š PDF ì´ë¯¸ì§€ ë¶„ì„:")
        print(f"      - ì „ì²´ ì´ë¯¸ì§€: {total_images}ê°œ")
        print(f"   ğŸ” í•„í„°ë§ í†µê³„:")
        print(f"      - ë°°ê²½ ì œì™¸: {filtered_background}ê°œ")
        print(f"      - ê°€ë¡œì„¸ë¡œë¹„: {filtered_aspect}ê°œ")
        print(f"      - ì‘ì€ ë©´ì : {filtered_area}ê°œ")
        print(f"      - ì‘ì€ í¬ê¸°: {filtered_size}ê°œ")
        print(f"   âœ… ìµœì¢… ì¶”ì¶œ: {len(metadata_list)}ê°œ ì´ë¯¸ì§€\n")
        
        return metadata_list


# 2. ê°œì„ ëœ í•˜ì´ë¸Œë¦¬ë“œ í•„í„° íŒŒì´í”„ë¼ì¸
class ImprovedHybridFilterPipeline:
    def __init__(self, auto_extract_keywords: bool = True):
        self.auto_extract = auto_extract_keywords
        
        self.UNIVERSAL_PATTERNS = [
            'í•™ìŠµ', 'í™œë™', 'ë¬¸ì œ', 'ì˜ˆì œ', 'ì—°ìŠµ',
            'ìƒê°', 'ì•Œì•„ë³´', 'ì‚´í´ë³´', 'ì •ë¦¬',
            'ëª©í‘œ', 'ê°œë…', 'ì›ë¦¬', 'ë²•ì¹™', 'ì •ì˜',
            'ë‹¨ì›', 'ì°¨ì‹œ',
            'ê·¸ë¦¼', 'ë„í‘œ', 'í‘œ', 'ì°¨íŠ¸', 'ê·¸ë˜í”„',
            'ì˜ˆì‹œ', 'ì‚¬ë¡€', 'ëª¨í˜•', 'êµ¬ì¡°'
        ]
        
        self.DECORATION_PATTERNS = [
            'ë¡œê³ ', 'logo', 'ì¶œì²˜', 'ì°¸ê³ ', 'ì•„ì´ì½˜', 'icon'
        ]
        
        self.document_keywords = []

    def extract_keywords_from_document(self, file_path: str):
        """ë¬¸ì„œì—ì„œ ìë™ìœ¼ë¡œ í‚¤ì›Œë“œ ì¶”ì¶œ"""
        if not self.auto_extract:
            return
        
        from pathlib import Path
        
        print("ğŸ“š ë¬¸ì„œ ë¶„ì„í•˜ì—¬ í‚¤ì›Œë“œ ìë™ ì¶”ì¶œ ì¤‘...")
        
        ext = Path(file_path).suffix.lower()
        all_text = []
        
        if ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
        
        elif ext == '.pdf':
            import pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            all_text.append(text)
            except Exception as e:
                print(f"   âš ï¸ PDF í…ìŠ¤íŠ¸ ì¶”ì¶œ ì‹¤íŒ¨, ë²”ìš© íŒ¨í„´ë§Œ ì‚¬ìš©")
                return
        
        else:
            print(f"   âš ï¸ ì§€ì›í•˜ì§€ ì•ŠëŠ” í˜•ì‹: {ext}")
            return
        
        full_text = "\n".join(all_text)[:5000]
        
        prompt = f"""
ë‹¤ìŒ ê°•ì˜ ìë£Œì—ì„œ **í•µì‹¬ í‚¤ì›Œë“œ 20ê°œ**ë¥¼ ì¶”ì¶œí•˜ì„¸ìš”.

# ë¬¸ì„œ ë‚´ìš©
{full_text}

# ì¡°ê±´
- ê°œë…ì–´, ì „ë¬¸ ìš©ì–´, ì£¼ì œì–´ë§Œ í¬í•¨
- JSON í˜•ì‹: {{"keywords": ["í‚¤ì›Œë“œ1", "í‚¤ì›Œë“œ2", ...]}}
"""
        
        try:
            response = model.generate_content(prompt)

            # âœ… í† í° ì¶”ì¶œ
            usage = response.usage_metadata
            in_t = usage.prompt_token_count
            out_t = usage.candidates_token_count
            cost = (in_t / 1_000_000 * 0.075) + (out_t / 1_000_000 * 0.30)
            
            print(f"ğŸ“Š [í‚¤ì›Œë“œ ì¶”ì¶œ] í† í°: {usage.total_token_count:,} (In: {in_t}/Out: {out_t}) / ë¹„ìš©: ${cost:.6f}")

            text = response.text.strip()
            
            if "```json" in text:
                text = text.split("```json")[1].split("```")[0].strip()
            elif "```" in text:
                text = text.split("```")[1].split("```")[0].strip()
            
            data = json.loads(text)
            self.document_keywords = data.get("keywords", [])
            
            print(f"   âœ… ì¶”ì¶œëœ í‚¤ì›Œë“œ: {', '.join(self.document_keywords[:10])}")
        
        except Exception as e:
            print(f"   âš ï¸ ìë™ ì¶”ì¶œ ì‹¤íŒ¨, ë²”ìš© íŒ¨í„´ë§Œ ì‚¬ìš©")
            self.document_keywords = []

    def step1_rule_check(self, meta: ImageMetadata):
        """ê·œì¹™ ê¸°ë°˜ 1ì°¨ í•„í„°"""
        context = f"{meta.slide_title} {meta.adjacent_text}".lower()
        
        has_deco = any(kw in context for kw in self.DECORATION_PATTERNS)
        is_corner = (meta.left < 1.0 and meta.top < 1.0) or (meta.left > 8.0 and meta.top < 1.0)
        
        if is_corner and meta.area_percentage < 5.0 and not any(kw in context for kw in self.UNIVERSAL_PATTERNS):
            return "EXCLUDE", "Static Decoration (Corner)"
        
        if has_deco and meta.area_percentage < 8.0:
            return "EXCLUDE", "Decorative element"
        
        has_universal = any(p in context for p in self.UNIVERSAL_PATTERNS)
        has_document_kw = any(kw in context for kw in self.document_keywords)
        
        if meta.area_percentage > 15.0 and (has_universal or has_document_kw):
            return "INCLUDE", f"Core content ({meta.area_percentage:.1f}% + pattern)"
        
        if has_document_kw and meta.area_percentage > 10.0:
            matched = [kw for kw in self.document_keywords if kw in context]
            return "INCLUDE", f"Document keyword: {', '.join(matched[:2])}"
        
        return "PENDING", "Requires AI Vision Check"

    def step2_gemini_check(self, meta: ImageMetadata, max_retries=3):
        """AI Visionìœ¼ë¡œ 2ì°¨ íŒë‹¨"""
        import time
        
        for attempt in range(max_retries):
            try:
                image_part = Part.from_data(data=meta.image_bytes, mime_type="image/png")
                
                keyword_list = ', '.join(list(self.document_keywords)[:15]) if self.document_keywords else "ì¼ë°˜ í•™ìŠµ ë‚´ìš©"
                
                prompt = f"""
ì´ ê°•ì˜ì˜ í•µì‹¬ ì£¼ì œ: {keyword_list}

ì´ ì´ë¯¸ì§€ê°€ ìœ„ ì£¼ì œë“¤ê³¼ ê´€ë ¨ìˆëŠ”ì§€ íŒë‹¨í•˜ì„¸ìš”.

ì£¼ë³€ í…ìŠ¤íŠ¸: "{meta.adjacent_text}"

íŒë‹¨:
- í•™ìŠµì— í•„ìš”í•œ í•µì‹¬ ìë£Œ â†’ KEEP + ì´ìœ 
- ì¥ì‹/ë¡œê³ /ë°°ê²½ â†’ DISCARD + ì´ìœ 

ì¶œë ¥: KEEP ë˜ëŠ” DISCARDë¡œ ì‹œì‘
"""
                response = model.generate_content([image_part, prompt])

                # âœ… í† í° ë° ë¹„ìš© ê³„ì‚° ì¶”ê°€
                # âœ… Gemini 2.5 Flash ê³µì‹ ë‹¨ê°€ ì ìš©
                input_tokens = response.usage_metadata.prompt_token_count
                output_tokens = response.usage_metadata.candidates_token_count
                total_tokens = response.usage_metadata.total_token_count
                cost = (input_tokens / 1_000_000 * 0.075) + (output_tokens / 1_000_000 * 0.30)

                return response.text.strip(), total_tokens, cost
                
            except Exception as e:
                error_msg = str(e)
                
                if "429" in error_msg or "Resource exhausted" in error_msg:
                    if attempt < max_retries - 1:
                        wait_time = (attempt + 1) * 3
                        print(f"      âš ï¸  Rate Limit, {wait_time}ì´ˆ ëŒ€ê¸°...")
                        time.sleep(wait_time)
                        continue
                    else:
                        return "DISCARD: API rate limit exceeded", 0, 0.0
                else:
                    return f"ERROR: {error_msg}", 0, 0.0
        
        return "DISCARD: Failed after all retries", 0, 0.0

    def run(self, source_path: str):
        """ì´ë¯¸ì§€ í•„í„°ë§ ì‹¤í–‰"""
        from pathlib import Path
        
        file_ext = Path(source_path).suffix.lower()
        print(f"\nğŸ” ë¶„ì„ ì‹œì‘: {os.path.basename(source_path)} ({file_ext})")
        
        if self.auto_extract:
            self.extract_keywords_from_document(source_path)
        
        extractor = UniversalImageExtractor()
        all_meta = extractor.extract(source_path)
        
        print("\n" + "="*120)
        print(f"{'Slide':<6} | {'Size':<6} | {'Filter':<12} | {'Result':<12} | {'Reason'}")
        print("-" * 120)

        final_core = []
        stats = {
            'total': len(all_meta),
            'rule_pass': 0,
            'rule_drop': 0,
            'ai_keep': 0,
            'ai_drop': 0,
        }
        
        for meta in all_meta:
            decision_type, s1_reason = self.step1_rule_check(meta)

            final_status = ""
            filter_stage = ""
            detail_reason = ""

            if decision_type == "INCLUDE":
                meta.is_core_content = True
                filter_stage = "1ì°¨ (Rule)"
                final_status = "âœ… PASS"
                detail_reason = s1_reason
                final_core.append(meta)
                stats['rule_pass'] += 1
                
            elif decision_type == "PENDING":
                filter_stage = "2ì°¨ (AI)"
                ai_res, tokens, cost = self.step2_gemini_check(meta)
                
                if ai_res.upper().startswith("KEEP"):
                    meta.is_core_content = True
                    final_status = "âœ… KEEP"
                    stats['ai_keep'] += 1
                    final_core.append(meta)
                else:
                    final_status = "âŒ DROP"
                    stats['ai_drop'] += 1
                    
                detail_reason = ai_res.replace('\n', ' ')
                
            else:
                filter_stage = "1ì°¨ (Rule)"
                final_status = "âŒ DROP"
                detail_reason = s1_reason
                stats['rule_drop'] += 1

            wrapped_reason = textwrap.wrap(detail_reason, width=70)
            print(f"{meta.slide_number:<6} | {meta.area_percentage:>5.1f}% | {filter_stage:<12} | {final_status:<12} | {wrapped_reason[0]}")
            for line in wrapped_reason[1:]:
                print(f"{'':<6} | {'':<6} | {'':<12} | {'':<12} | {line}")
            print("-" * 120)

        print("\n" + "="*120)
        print("ğŸ“Š ìµœì¢… ê²°ê³¼")
        print("="*120)
        
        print(f"\nì´ ì´ë¯¸ì§€: {stats['total']}ê°œ")
        print(f"\n[1ì°¨ í•„í„° - ê·œì¹™ ê¸°ë°˜]")
        print(f"  âœ… í†µê³¼: {stats['rule_pass']}ê°œ")
        print(f"  âŒ ì œì™¸: {stats['rule_drop']}ê°œ")
        print(f"  âš ï¸  2ì°¨ ì´ë™: {stats['ai_keep'] + stats['ai_drop']}ê°œ")
        
        print(f"\n[2ì°¨ í•„í„° - AI íŒë‹¨]")
        print(f"  âœ… í†µê³¼: {stats['ai_keep']}ê°œ")
        print(f"  âŒ ì œì™¸: {stats['ai_drop']}ê°œ")
        
        total_keep = stats['rule_pass'] + stats['ai_keep']
        total_drop = stats['rule_drop'] + stats['ai_drop']
        
        print(f"\n{'='*120}")
        print(f"ğŸ’ ìµœì¢… í•µì‹¬ ì´ë¯¸ì§€: {total_keep}ê°œ (1ì°¨: {stats['rule_pass']}ê°œ + 2ì°¨: {stats['ai_keep']}ê°œ)")
        print(f"ğŸ—‘ï¸  ì œì™¸ëœ ì´ë¯¸ì§€: {total_drop}ê°œ")
        if stats['total'] > 0:
            print(f"ğŸ’° Vision API ì‚¬ìš©: {stats['ai_keep'] + stats['ai_drop']}íšŒ ({(stats['ai_keep'] + stats['ai_drop'])/stats['total']*100:.1f}%)")
        print(f"{'='*120}\n")
        
        return final_core


if __name__ == "__main__":
    import sys
    
    print("\n" + "="*120)
    print("ğŸ¯ Improved Hybrid Filter V2 - ì´ë¯¸ì§€ í•„í„°ë§")
    print("="*120)
    
    if len(sys.argv) > 1:
        source_file = sys.argv[1]
        
        if not os.path.exists(source_file):
            print(f"\nâŒ íŒŒì¼ì„ ì°¾ì„ ìˆ˜ ì—†ìŠµë‹ˆë‹¤: {source_file}")
            sys.exit(1)
        
        auto_extract = True
        if len(sys.argv) > 2 and sys.argv[2] in ['--no-auto', '-n']:
            auto_extract = False
            print("\nâš ï¸  ìë™ í‚¤ì›Œë“œ ì¶”ì¶œ ë¹„í™œì„±í™”")
        else:
            print("\nâœ… ìë™ í‚¤ì›Œë“œ ì¶”ì¶œ í™œì„±í™”")
        
        try:
            pipeline = ImprovedHybridFilterPipeline(auto_extract_keywords=auto_extract)
            core_images = pipeline.run(source_file)
            
            print(f"\n{'='*120}")
            print(f"âœ… ì™„ë£Œ! í•µì‹¬ ì´ë¯¸ì§€: {len(core_images)}ê°œ")
            print(f"{'='*120}\n")
            
        except Exception as e:
            print(f"\nâŒ ì—ëŸ¬ ë°œìƒ: {e}")
            import traceback
            traceback.print_exc()
            sys.exit(1)
    
    else:
        print("\nì‚¬ìš©ë²•:")
        print("  python improved_hybrid_filter_v2.py <íŒŒì¼ê²½ë¡œ>")
        print("\nì˜ˆì‹œ:")
        print("  python improved_hybrid_filter_v2.py ì¤‘ë“±êµ­ì–´1.pdf")
        print("\nâœ… V2 ê°œì„ ì‚¬í•­:")
        print("  - get_images() ë°©ì‹ìœ¼ë¡œ ëª¨ë“  ì´ë¯¸ì§€ ê°ì§€")
        print("  - ë§Œí™” ì½˜í…ì¸  ì •ìƒ ì¸ì‹")
        print("  - ë°°ê²½ ì´ë¯¸ì§€ ìë™ ì œì™¸")
        print("="*120 + "\n")